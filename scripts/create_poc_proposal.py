# -*- coding: utf-8 -*-
"""
SENTIO AI - PoC 병원 제안서 PPTX 생성 스크립트 v2
sentio-1 디자인 테마 기반: Pretendard, 미니멀 화이트, #272525/#C1001B
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

# ── 경로 ──
IMG_DIR = r"C:\Users\dbals\VibeCoding\Care-guard\docs\01_presentation\gen_images"
DASH_DIR = r"C:\Users\dbals\VibeCoding\Care-guard\docs\01_presentation"
LOGO_PATH = r"C:\Users\dbals\VibeCoding\Care-guard\docs\01_presentation\sentio_logo_sm.png"
OUTPUT = r"C:\Users\dbals\VibeCoding\Care-guard\SENTIO_PoC_제안서.pptx"
OUTPUT_ONEDRIVE = r"C:\Users\dbals\OneDrive\문서\카카오톡 받은 파일\SENTIO_PoC_제안서.pptx"

# ── sentio-1 테마 색상 ──
C_TEXT = RGBColor(0x27, 0x25, 0x25)       # 본문
C_BLACK = RGBColor(0x00, 0x00, 0x00)      # 제목
C_RED = RGBColor(0xC1, 0x00, 0x1B)        # 강조
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xF7, 0xF7, 0xF7)     # 밝은 배경
C_BORDER = RGBColor(0xE5, 0xE5, 0xE5)    # 테두리
C_BLUE = RGBColor(0x25, 0x63, 0xEB)      # SENTIO 브랜드 블루
C_GREEN = RGBColor(0x05, 0x96, 0x69)     # 안전 녹색
C_YELLOW = RGBColor(0xD9, 0x77, 0x06)    # 주의 노란
C_GRAY = RGBColor(0x9C, 0x9C, 0x9C)      # 보조 텍스트

# ── 폰트 크기 (sentio-1 기준 EMU → Pt 변환) ──
SZ_CATEGORY = Pt(13)    # 165100 EMU → 카테고리 태그
SZ_TITLE = Pt(33.5)     # 425450 EMU → 큰 제목
SZ_SUBTITLE = Pt(16.5)  # 209550 EMU → 소제목/라벨
SZ_BODY = Pt(16)         # 203200 EMU → 본문
SZ_STAT = Pt(40)         # 514350 EMU → 큰 숫자
SZ_SMALL = Pt(14)        # 소형 본문
SZ_TINY = Pt(12)         # 작은 보조

FONT = "Pretendard"
FONT_B = "Pretendard Bold"

# ── 프레젠테이션 ──
prs = Presentation()
prs.slide_width = Emu(14630400)  # sentio-1과 동일
prs.slide_height = Emu(8229600)
SW = 14630400
SH = 8229600


# ══════════════════════════════════════════
# 유틸리티
# ══════════════════════════════════════════

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def tb(slide, l, t, w, h, text, sz=SZ_BODY, color=C_TEXT, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    """텍스트 박스 (한국어 어절 줄바꿈 지원)"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    # 한국어 어절 줄바꿈: Latin/EA line break 설정
    for attr in ['latinLnBrk', 'eaLnBrk']:
        tf._txBody.bodyPr.set(attr, '1')
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = sz
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def para(tf, text, sz=SZ_BODY, color=C_TEXT, bold=False, align=PP_ALIGN.LEFT, font=FONT,
         sp_before=Pt(6), sp_after=Pt(2)):
    """텍스트 프레임에 단락 추가"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = sz
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_before = sp_before
    p.space_after = sp_after
    return p

def rect(slide, l, t, w, h, fill=None, border=None):
    """사각형"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill=None, border=None):
    """둥근 사각형"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def circle(slide, l, t, size, fill=C_TEXT):
    """원"""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def dot(slide, l, t, color=C_TEXT):
    """작은 불릿 점"""
    return circle(slide, l, t, Cm(0.25), color)

def img(slide, path, l, t, w=None, h=None, width=None, height=None):
    """이미지 (안전, 비율 유지 - width 또는 height 하나만 지정)"""
    if not os.path.exists(path):
        return None
    _w = w or width
    _h = h or height
    # 비율 유지: 하나만 지정
    if _w and _h:
        # 둘 다 지정되면 width 기준으로 비율 유지
        return slide.shapes.add_picture(path, l, t, width=_w)
    elif _w:
        return slide.shapes.add_picture(path, l, t, width=_w)
    elif _h:
        return slide.shapes.add_picture(path, l, t, height=_h)
    else:
        return slide.shapes.add_picture(path, l, t)

def add_logo(slide):
    """좌하단 SENTIO 로고 추가"""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Cm(1), Emu(SH - Cm(1.5)), height=Cm(0.9))

def slide_header(slide, category, title, cat_top=None, title_top=None):
    """sentio-1 스타일 헤더: 작은 카테고리 + 큰 타이틀 + 로고"""
    ct = cat_top or Cm(2.2)
    tt = title_top or Cm(2.8)
    tb(slide, Cm(2.2), ct, Cm(20), Cm(0.7), category,
       sz=SZ_CATEGORY, color=C_TEXT, font=FONT)
    tb(slide, Cm(2), tt, Cm(28), Cm(1.5), title,
       sz=SZ_TITLE, color=C_BLACK, bold=True, font=FONT_B)
    # 우상단 SENTIO 로고
    add_logo(slide)

def slide_num(slide, n, total=20):
    """슬라이드 번호"""
    tb(slide, Emu(SW - Cm(3)), Emu(SH - Cm(1)), Cm(2.5), Cm(0.5),
       f"{n}", sz=SZ_TINY, color=C_GRAY, align=PP_ALIGN.RIGHT)

def bottom_note(slide, text, color=C_TEXT):
    """하단 설명"""
    # 구분선
    rect(slide, Cm(2), Emu(SH - Cm(2.5)), Emu(SW - Cm(4)), Pt(0.5), C_BORDER)
    tb(slide, Cm(2), Emu(SH - Cm(2.2)), Emu(SW - Cm(4)), Cm(1.5), text,
       sz=SZ_BODY, color=color, font=FONT)


# ══════════════════════════════════════════
# 슬라이드 1: 표지
# ══════════════════════════════════════════
def s01():
    slide = blank()
    # 좌측 이미지 (반쪽 커버) - 너비 기준 비율 유지 (1:1 이미지)
    img(slide, os.path.join(IMG_DIR, "slide3_ai_dashboard.png"),
        0, 0, width=Emu(SW // 2 - Cm(1)))

    # 우측 텍스트 영역
    x = Emu(SW // 2 + Cm(1))
    tb(slide, x, Cm(3), Cm(14), Cm(1.5), "SENTIO AI",
       sz=Pt(46), color=C_BLACK, bold=True, font=FONT_B)
    tb(slide, x, Cm(5.5), Cm(14), Cm(1.5), "요양병원 낙상 감지 시스템",
       sz=Pt(28), color=C_BLACK, bold=True, font=FONT_B)

    # PoC 제안서 태그
    tag = rrect(slide, x, Cm(8), Cm(5), Cm(1), fill=C_RED)
    tf = tag.text_frame
    tf.paragraphs[0].text = "PoC 제안서"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = FONT_B
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 핵심 문구
    tb(slide, x, Cm(10), Cm(16), Cm(0.8),
       "20건의 낙상 중 19건을 자동 감지  |  424개 실제 영상으로 검증 완료",
       sz=SZ_BODY, color=C_TEXT, font=FONT)
    tb(slide, x, Cm(11.5), Cm(10), Cm(0.5),
       "2026년 3월", sz=SZ_TINY, color=C_GRAY)
    # 표지에도 로고 추가
    add_logo(slide)


# ══════════════════════════════════════════
# 슬라이드 2: Executive Summary
# ══════════════════════════════════════════
def s02():
    slide = blank()
    slide_header(slide, "Executive Summary", "의사결정자를 위한 핵심 요약")
    slide_num(slide, 2)

    y = Cm(5.5)
    labels = [
        ("문제", C_RED,
         "요양병원 낙상 사고 — 연간 28~35% 발생, 환자안전사고 2위"),
        ("솔루션", C_BLUE,
         "기존 CCTV + AI로 낙상 즉시 감지 — 추가 장비 없이 도입 가능"),
        ("제안", C_GREEN,
         "2~4주 무료 PoC — 귀 병원 환경에서 직접 검증 후 결정"),
    ]
    for label, color, msg in labels:
        # 불릿 도트
        dot(slide, Cm(2.2), y + Cm(0.2), color)
        tb(slide, Cm(3), y, Cm(3), Cm(0.7), label,
           sz=SZ_SUBTITLE, color=color, bold=True, font=FONT_B)
        tb(slide, Cm(6), y, Cm(26), Cm(0.7), msg,
           sz=SZ_BODY, color=C_TEXT)
        y += Cm(1.8)

    # 제안 개요 박스
    y += Cm(0.5)
    box = rrect(slide, Cm(2), y, Emu(SW - Cm(4)), Cm(5.5), fill=C_LIGHT, border=C_BORDER)
    tb(slide, Cm(3), y + Cm(0.4), Cm(10), Cm(0.7), "제안 개요",
       sz=SZ_SUBTITLE, color=C_BLACK, bold=True, font=FONT_B)

    items = [
        ("대상", "요양병원 · 요양시설 (기존 CCTV 보유 기관)"),
        ("내용", "AI 기반 실시간 낙상 감지 시스템 'SENTIO' PoC 운영"),
        ("기간", "2~4주 (설치 1일 + 운영 + 결과 리포트)"),
        ("비용", "PoC 기간 전액 무료 — 성과 확인 후 도입 결정"),
    ]
    iy = y + Cm(1.5)
    for lbl, val in items:
        tb(slide, Cm(3.5), iy, Cm(2.5), Cm(0.6), lbl,
           sz=SZ_SMALL, color=C_RED, bold=True, font=FONT_B)
        tb(slide, Cm(6.5), iy, Cm(26), Cm(0.6), val,
           sz=SZ_SMALL, color=C_TEXT)
        iy += Cm(0.9)


# ══════════════════════════════════════════
# 슬라이드 3: 문제 인식
# ══════════════════════════════════════════
def s03():
    slide = blank()
    slide_header(slide, "문제 인식", "낙상 사고, 얼마나 심각한가?")
    slide_num(slide, 3)

    stats = [
        ("28~35%", "연간 낙상률", "65세 이상 노인 기준"),
        ("2위", "의료사고 원인", "전체 의료사고 중 낙상 순위"),
        ("47%", "대기실·복도 낙상", "전체 낙상 중 비율"),
        ("20~30명", "간호사 1인당 담당", "동시 모니터링 한계"),
    ]
    x = Cm(2)
    for num, label, sub in stats:
        tb(slide, x, Cm(5.5), Cm(7.5), Cm(2), num,
           sz=Pt(52), color=C_TEXT, bold=True, font=FONT_B, align=PP_ALIGN.LEFT)
        tb(slide, x + Cm(0.3), Cm(8), Cm(7), Cm(0.7), label,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, x + Cm(0.3), Cm(8.8), Cm(7), Cm(0.6), sub,
           sz=SZ_BODY, color=C_TEXT)
        x += Cm(8.5)

    bottom_note(slide,
        "낙상 20건 중 19건을 자동 감지하면서, 잘못된 알림으로 인한 업무 방해를 최소화합니다.")


# ══════════════════════════════════════════
# 슬라이드 4: 현장의 어려움
# ══════════════════════════════════════════
def s04():
    slide = blank()
    # 우측 이미지 - 너비 기준 비율 유지 (1:1 이미지)
    img(slide, os.path.join(IMG_DIR, "scene1_hospital_fall.png"),
        Emu(SW // 2 + Cm(1)), 0, width=Emu(SW // 2 - Cm(1)))

    slide_header(slide, "현장의 어려움", "24시간 감시는 현실적으로 불가능합니다")
    slide_num(slide, 4)

    items = [
        ("물리적 한계", "간호사 1인이 대기실·복도를 24시간 지켜보는 것은 구조적으로 불가능합니다."),
        ("연쇄 피해", "낙상 → 골절 → 입원 연장 → 의료비 증가 → 법적 책임으로 이어지는 악순환"),
        ("골든타임", "발견이 늦을수록 피해는 커집니다. 신속한 대응이 핵심입니다."),
    ]
    y = Cm(6)
    for title, desc in items:
        dot(slide, Cm(2.2), y + Cm(0.15), C_RED)
        tb(slide, Cm(3), y, Cm(13), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, Cm(3), y + Cm(0.9), Cm(14), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        y += Cm(3.5)


# ══════════════════════════════════════════
# 슬라이드 5: 솔루션 개요
# ══════════════════════════════════════════
def s05():
    slide = blank()
    slide_header(slide, "솔루션 소개", "SENTIO는 무엇인가?")
    slide_num(slide, 5)

    tb(slide, Cm(3), Cm(4.8), Emu(SW - Cm(6)), Cm(0.8),
       "기존 CCTV에 AI를 연결하면, 환자가 넘어지는 순간 간호사에게 즉시 알려줍니다.",
       sz=SZ_BODY, color=C_TEXT)

    # 이미지 아이콘 대신 원형 + 텍스트 (sentio-1 스타일)
    features = [
        ("기존 CCTV 그대로 활용", "새로운 장비 구매 없이 즉시 도입"),
        ("AI 실시간 감지", "넘어지는 순간 자동 알림 → 빠른 대응"),
        ("프라이버시 보호", "얼굴·신체 노출 없이 관절점만 분석"),
        ("여러 환자 동시 추적", "대기실 다수 인원 개별 모니터링"),
    ]
    # 2x2 그리드 - 더 넓게 분포
    positions = [
        (Cm(2), Cm(7)), (Cm(18), Cm(7)),
        (Cm(2), Cm(10)), (Cm(18), Cm(10)),
    ]
    for (title, desc), (x, y) in zip(features, positions):
        dot(slide, x, y + Cm(0.15), C_RED)
        tb(slide, x + Cm(0.7), y, Cm(12), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, x + Cm(0.7), y + Cm(0.7), Cm(14), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)

    # CCTV 이미지 - 적절한 크기
    img(slide, os.path.join(IMG_DIR, "scene7_cctv.png"),
        Cm(26), Cm(14), height=Cm(7))


# ══════════════════════════════════════════
# 슬라이드 6: 경쟁력 비교
# ══════════════════════════════════════════
def s06():
    slide = blank()
    slide_header(slide, "비교 분석", "SENTIO vs 기존 방식")
    slide_num(slide, 6)

    # 테이블 데이터
    headers = ["항목", "기존 CCTV", "기존 센서 솔루션", "SENTIO"]
    col_x = [Cm(2.5), Cm(10.5), Cm(18), Cm(25.5)]
    col_w = [Cm(7.5), Cm(7), Cm(7), Cm(7)]

    # 헤더
    y = Cm(5.5)
    for cx, cw, h in zip(col_x, col_w, headers):
        is_sentio = (h == "SENTIO")
        tb(slide, cx, y, cw, Cm(0.7), h,
           sz=SZ_BODY, color=C_RED if is_sentio else C_TEXT, bold=True, font=FONT_B)
    # 구분선
    rect(slide, Cm(2.5), y + Cm(0.8), Emu(SW - Cm(5)), Pt(1), C_TEXT)

    rows = [
        ("낙상 감지", "사람이 직접 확인", "별도 센서 필요", "AI 자동 감지"),
        ("추가 장비", "—", "1,000만원+", "불필요"),
        ("프라이버시", "원본 저장", "해당 없음", "관절점만 처리"),
        ("다중 추적", "불가", "제한적", "동시 추적 가능"),
        ("감지 속도", "발견 시까지", "센서 지연", "3초 이내"),
        ("월 비용", "—", "50만원+", "30만원~"),
    ]
    y = Cm(7)
    for row in rows:
        # 얇은 구분선
        rect(slide, Cm(2.5), y - Cm(0.1), Emu(SW - Cm(5)), Pt(0.3), C_BORDER)
        for cx, cw, val in zip(col_x, col_w, row):
            is_sentio_col = (cx == col_x[3])
            tb(slide, cx, y, cw, Cm(0.8), val,
               sz=SZ_BODY, color=C_RED if is_sentio_col else C_TEXT,
               bold=is_sentio_col, font=FONT_B if is_sentio_col else FONT)
        y += Cm(1.8)


# ══════════════════════════════════════════
# 슬라이드 7: 작동 원리
# ══════════════════════════════════════════
def s07():
    slide = blank()
    slide_header(slide, "작동 원리", "어떻게 동작하나?", cat_top=Cm(0.8), title_top=Cm(1.4))
    slide_num(slide, 7)

    # 작동 원리 이미지 (기존 PPTX에서 추출한 것 또는 gen_images)
    # 4단계 플로우
    steps = ["CCTV 입력", "AI 사람 인식", "낙상 판단", "즉시 알림"]
    x = Cm(2)
    y = Cm(4)
    for i, step in enumerate(steps):
        # 둥근 박스
        box = rrect(slide, x, y, Cm(6.5), Cm(1.2), fill=C_LIGHT, border=C_BORDER)
        tf = box.text_frame
        tf.paragraphs[0].text = step
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = C_WHITE if i == 3 else C_TEXT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_B
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if i == 3:
            box.fill.solid()
            box.fill.fore_color.rgb = C_RED

        # 화살표
        if i < 3:
            tb(slide, x + Cm(7), y + Cm(0.1), Cm(1.5), Cm(1), "→",
               sz=Pt(24), color=C_GRAY, align=PP_ALIGN.CENTER)
        x += Cm(8.5)

    # 하단 설명
    tb(slide, Cm(2), Cm(6), Emu(SW - Cm(4)), Cm(0.7),
       "기존 CCTV 영상을 실시간으로 분석하여, 낙상 감지 즉시 소리·화면 경고·스마트폰 푸시 알림을 동시에 발송합니다.",
       sz=SZ_BODY, color=C_TEXT)

    # 대시보드 이미지 - 높이 기준 비율 유지 (overflow 방지)
    img(slide, os.path.join(DASH_DIR, "dashboard_normal.png"),
        Cm(2), Cm(7.5), height=Cm(13.5))


# ══════════════════════════════════════════
# 슬라이드 8: AI 분석 5단계
# ══════════════════════════════════════════
def s08():
    slide = blank()
    slide_header(slide, "AI 분석 과정", "AI가 하는 5가지 일")
    slide_num(slide, 8)

    stages = [
        ("01", "사람 찾기", "영상에서 사람을 자동으로 탐지합니다."),
        ("02", "관절 추적", "어깨·엉덩이·무릎 등 17개 주요 관절 움직임을 추적합니다."),
        ("03", "패턴 분석", "66만 건 학습 데이터로 정상/비정상 움직임을 구분합니다."),
        ("04", "위험 판정", "낙상 확률 계산 — 임계값 초과 시 위험으로 판정합니다."),
        ("05", "즉시 알림", "2단계 확인 후 3초 이내에 소리·화면·푸시 알림을 발송합니다."),
    ]
    y = Cm(5)
    for num, title, desc in stages:
        # 번호
        tb(slide, Cm(2.5), y, Cm(2), Cm(0.8), num,
           sz=Pt(20), color=C_RED, bold=True, font=FONT_B)
        # 제목
        tb(slide, Cm(5), y, Cm(6), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        # 설명
        tb(slide, Cm(12), y, Cm(20), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        # 구분선
        if num != "05":
            rect(slide, Cm(2.5), y + Cm(0.9), Emu(SW - Cm(5)), Pt(0.3), C_BORDER)
        y += Cm(1.2)

    # AI 포즈 이미지 - 우하단, 텍스트와 겹치지 않게
    img(slide, os.path.join(IMG_DIR, "slide5_ai_pose_analysis.png"),
        Cm(25), Cm(12), height=Cm(9))


# ══════════════════════════════════════════
# 슬라이드 9: 3단계 알림
# ══════════════════════════════════════════
def s09():
    slide = blank()
    slide_header(slide, "알림 체계", "3단계 알림 시스템")
    slide_num(slide, 9)

    levels = [
        ("1단계 — 정상", C_GREEN, "이상 없음. 모니터링만 지속합니다.",
         "화면 표시만"),
        ("2단계 — 주의", C_YELLOW, "이상 징후 감지. 화면 경고가 표시됩니다.",
         "화면 경고"),
        ("3단계 — 위험", C_RED, "낙상 감지. 즉각 대응이 필요합니다.",
         "소리 + 화면 + 푸시 + 음성 안내 동시 발령"),
    ]

    x = Cm(2)
    for name, color, desc, action in levels:
        # 상태 아이콘 (원)
        c = circle(slide, x, Cm(5.5), Cm(1), color)
        # 이름
        tb(slide, x + Cm(1.5), Cm(5.3), Cm(9), Cm(0.7), name,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        # 설명
        tb(slide, x + Cm(1.5), Cm(6.1), Cm(9), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        # 액션
        tb(slide, x + Cm(1.5), Cm(6.9), Cm(9), Cm(0.7), action,
           sz=SZ_SMALL, color=color, bold=True, font=FONT_B)
        x += Cm(11)

    # 알림 대시보드 - 높이 기준 비율 유지 (overflow 방지)
    img(slide, os.path.join(DASH_DIR, "dashboard_alert.png"),
        Cm(2), Cm(9), height=Cm(12))


# ══════════════════════════════════════════
# 슬라이드 10: 프라이버시
# ══════════════════════════════════════════
def s10():
    slide = blank()
    # 우측 이미지 - 너비 기준 비율 유지 (1:1 이미지)
    img(slide, os.path.join(IMG_DIR, "scene7_privacy.png"),
        Emu(SW // 2 + Cm(1)), 0, width=Emu(SW // 2 - Cm(1)))

    slide_header(slide, "프라이버시", "환자 프라이버시, 어떻게 보호하나?")
    slide_num(slide, 10)

    # 3가지 모드 - sentio-1 스타일 (리스트형)
    modes = [
        ("스켈레톤 (기본)", "검정 배경 + 관절점만 표시. 사람의 외형이 전혀 보이지 않음."),
        ("블러 모드", "사람 영역을 자동 블러 처리. 움직임은 확인, 신원은 보호."),
        ("원본 모드 (관리자용)", "관리자만 접근 가능. 사고 확인 시에만 제한적 열람."),
    ]
    y = Cm(6)
    for title, desc in modes:
        dot(slide, Cm(2.2), y + Cm(0.15), C_RED)
        tb(slide, Cm(3), y, Cm(13), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, Cm(3), y + Cm(0.9), Cm(14), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        y += Cm(3)

    # 핵심 원칙
    tb(slide, Cm(2.2), Emu(SH - Cm(2.5)), Cm(16), Cm(0.7),
       "핵심: 원본 영상 미저장  ·  AI 분석 후 즉시 폐기  ·  개인정보보호법 준수",
       sz=SZ_SMALL, color=C_RED, bold=True, font=FONT_B)


# ══════════════════════════════════════════
# 슬라이드 11: 검증된 성능
# ══════════════════════════════════════════
def s11():
    slide = blank()
    slide_header(slide, "성능 검증", "424개 실제 영상으로 검증한 정확도")
    slide_num(slide, 11)

    metrics = [
        ("95.2%", "감지율 (Recall)", "실제 낙상 20건 중 19건 자동 감지"),
        ("94.6%", "정밀도 (Precision)", "알림이 울렸을 때 95% 확률로 진짜 낙상"),
        ("94.9", "F1 Score", "감지율·정밀도 균형 지표"),
    ]
    x = Cm(2)
    for num, label, desc in metrics:
        tb(slide, x, Cm(5), Cm(10), Cm(2), num,
           sz=Pt(52), color=C_TEXT, bold=True, font=FONT_B, align=PP_ALIGN.LEFT)
        tb(slide, x + Cm(0.3), Cm(7.2), Cm(9), Cm(0.7), label,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, x + Cm(0.3), Cm(8), Cm(9), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        x += Cm(11.5)

    bottom_note(slide,
        "낙상 20건 중 19건을 자동 감지하면서, 잘못된 알림으로 인한 업무 방해를 최소화합니다. "
        "간호사의 부담을 줄이면서도 환자 안전을 지키는 균형 잡힌 성능입니다.")


# ══════════════════════════════════════════
# 슬라이드 12: 유형별 감지율
# ══════════════════════════════════════════
def s12():
    slide = blank()
    slide_header(slide, "유형별 성능", "낙상 유형별 감지율")
    slide_num(slide, 12)

    types = [
        ("뒤로 넘어짐 — 99%", "가장 위험한 유형. 거의 100% 감지합니다."),
        ("앞으로 넘어짐 — 97%", "앞으로 쓰러지는 패턴을 정확히 인식합니다."),
        ("옆으로 넘어짐 — 91%", "옆으로 쓰러지는 패턴도 높은 정확도로 감지합니다."),
    ]
    y = Cm(5.5)
    for title, desc in types:
        dot(slide, Cm(2.2), y + Cm(0.15), C_RED)
        tb(slide, Cm(3), y, Cm(28), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, Cm(3), y + Cm(0.9), Cm(28), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        y += Cm(2.5)

    # 경쟁 비교 요약
    y += Cm(1)
    rect(slide, Cm(2), y, Emu(SW - Cm(4)), Pt(0.5), C_BORDER)
    y += Cm(0.3)
    tb(slide, Cm(2.5), y, Cm(28), Cm(0.7), "경쟁력 비교",
       sz=SZ_SUBTITLE, color=C_BLACK, bold=True, font=FONT_B)

    comp_headers = ["구분", "SENTIO", "A 연구 (2024)", "B 연구 (2024)", "기존 센서"]
    cx = [Cm(2.5), Cm(9), Cm(16), Cm(22.5), Cm(29)]
    cw = Cm(6)
    y += Cm(0.8)
    for c, h in zip(cx, comp_headers):
        tb(slide, c, y, cw, Cm(0.6), h,
           sz=SZ_SMALL, color=C_RED if h == "SENTIO" else C_TEXT, bold=True, font=FONT_B)
    rect(slide, Cm(2.5), y + Cm(0.6), Emu(SW - Cm(5)), Pt(0.5), C_BORDER)

    comp_rows = [
        ("검증 영상", "424건", "60건", "30건", "소규모"),
        ("감지율", "95.2%", "92~95%", "93~95%", "60~80%"),
    ]
    for row in comp_rows:
        y += Cm(0.7)
        for c, val in zip(cx, row):
            is_s = (c == cx[1])
            tb(slide, c, y, cw, Cm(0.6), val,
               sz=SZ_SMALL, color=C_RED if is_s else C_TEXT,
               bold=is_s, font=FONT_B if is_s else FONT)


# ══════════════════════════════════════════
# 슬라이드 13: 신뢰성 근거
# ══════════════════════════════════════════
def s13():
    slide = blank()
    # 우측 이미지 - 너비 기준 비율 유지 (1:1 이미지)
    img(slide, os.path.join(IMG_DIR, "scene3_ai_technology.png"),
        Emu(SW // 2 + Cm(1)), 0, width=Emu(SW // 2 - Cm(1)))

    slide_header(slide, "신뢰성 근거", "왜 SENTIO의 수치를 신뢰할 수 있나?")
    slide_num(slide, 13)

    items = [
        ("대규모 검증", "424개 실제 영상 — 타 연구 대비 7~14배 많은 검증 데이터"),
        ("공인 데이터 사용", "AI Hub(한국지능정보사회진흥원) 공인 낙상 데이터셋 활용"),
        ("대용량 학습", "66만 건 이상의 움직임 샘플로 AI 훈련 완료"),
        ("투명한 성능 공개", "감지율·오탐율·미탐율 모두 공개 — 숨기는 수치 없음"),
    ]
    y = Cm(6)
    for title, desc in items:
        dot(slide, Cm(2.2), y + Cm(0.15), C_RED)
        tb(slide, Cm(3), y, Cm(14), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, Cm(3), y + Cm(0.9), Cm(15), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        y += Cm(2.8)


# ══════════════════════════════════════════
# 슬라이드 14: 주요 기능
# ══════════════════════════════════════════
def s14():
    slide = blank()
    slide_header(slide, "주요 기능", "주요 기능 한눈에")
    slide_num(slide, 14)

    features = [
        ("실시간 대시보드", "모든 카메라를 한 화면에서 통합 모니터링"),
        ("한국어 음성 경고", "TTS로 '낙상이 감지되었습니다' 즉시 안내"),
        ("스마트폰 알림", "간호사 스마트폰으로 실시간 푸시 알림"),
        ("SafeZone 설정", "병실·복도 등 영역별 맞춤 감시 설정"),
        ("EMR 연동", "기존 의료정보시스템과 연동 가능 (Pro+)"),
        ("리포트 생성", "낙상 이력·통계 리포트 자동 생성"),
    ]
    # 3x2 좌/우
    y = Cm(5)
    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Cm(2) if col == 0 else Cm(18)
        yy = y + row * Cm(2.2)
        dot(slide, x, yy + Cm(0.15), C_RED)
        tb(slide, x + Cm(0.7), yy, Cm(14), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, x + Cm(0.7), yy + Cm(0.7), Cm(14), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)

    # 앱 스크린샷 - 하단 중앙에 배치하여 텍스트와 겹치지 않게
    img(slide, os.path.join(IMG_DIR, "slide5_actual_app.png"),
        Cm(12), Cm(13), height=Cm(8))


# ══════════════════════════════════════════
# 슬라이드 15: PoC 제안 개요 ★★★
# ══════════════════════════════════════════
def s15():
    slide = blank()
    slide_header(slide, "PoC 제안", "귀 병원 환경에서 직접 검증하세요")
    slide_num(slide, 15)

    # 강조: 비용 부담 없이
    tb(slide, Cm(2), Cm(4.8), Cm(28), Cm(0.8),
       "비용 부담 없이, 2~4주간 실제 운영 환경에서 SENTIO의 성능을 직접 확인하실 수 있습니다.",
       sz=SZ_BODY, color=C_TEXT)

    # 좌측: PoC 개요
    items = [
        ("목적", "SENTIO의 낙상 감지 성능을 귀 병원의 실제 환경에서 검증"),
        ("기간", "2~4주 (설치 1일 포함)"),
        ("범위", "카메라 2~4대 (병실·복도 등 핵심 구역)"),
        ("비용", "PoC 기간 전액 무료"),
    ]
    y = Cm(7)
    for lbl, val in items:
        tb(slide, Cm(2.5), y, Cm(2.5), Cm(0.6), lbl,
           sz=SZ_SMALL, color=C_RED, bold=True, font=FONT_B)
        tb(slide, Cm(5.5), y, Cm(12), Cm(0.6), val,
           sz=SZ_SMALL, color=C_TEXT)
        y += Cm(1.5)

    # 우측: 병원 준비사항 / SENTIO 제공
    rx = Cm(19)
    tb(slide, rx, Cm(6.5), Cm(13), Cm(0.7), "병원 측 준비사항",
       sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
    preps = [
        "기존 CCTV (IP 카메라, RTSP 지원)",
        "네트워크 연결 (유선 LAN 권장)",
        "모니터링용 PC 또는 태블릿 1대",
        "담당자 지정 (설치 협조)",
    ]
    py = Cm(7.5)
    for item in preps:
        tb(slide, rx + Cm(0.5), py, Cm(12), Cm(0.6), f"·  {item}",
           sz=SZ_SMALL, color=C_TEXT)
        py += Cm(0.9)

    tb(slide, rx, Cm(11.5), Cm(13), Cm(0.7), "SENTIO 제공사항",
       sz=SZ_SUBTITLE, color=C_RED, bold=True, font=FONT_B)
    provides = ["AI 분석 서버 설치 및 설정", "시스템 교육 (30분)", "PoC 결과 리포트"]
    sy = Cm(12.5)
    for item in provides:
        tb(slide, rx + Cm(0.5), sy, Cm(12), Cm(0.6), f"·  {item}",
           sz=SZ_SMALL, color=C_TEXT)
        sy += Cm(0.9)


# ══════════════════════════════════════════
# 슬라이드 16: PoC 진행 절차 ★★★
# ══════════════════════════════════════════
def s16():
    slide = blank()
    slide_header(slide, "도입 프로세스", "5단계 도입 프로세스")
    slide_num(slide, 16)

    steps = [
        ("문의 접수", "병원 환경 파악, 무료"),
        ("환경 분석", "CCTV·네트워크 호환성 확인, 무료"),
        ("무료 PoC", "2~4주 실제 운영 테스트"),
        ("결과 리포트", "감지 정확도·알림 통계 제공"),
        ("계약·운영", "정식 도입·교육·기술지원"),
    ]

    # 가로 타임라인 (sentio-1 slide 18 스타일)
    x = Cm(2)
    y = Cm(7)
    for i, (title, desc) in enumerate(steps):
        # 원형 번호
        c = circle(slide, x + Cm(2.5), y, Cm(1.5), C_RED if i == 2 else C_TEXT)
        tf = c.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = C_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_B
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 제목
        tb(slide, x, y + Cm(2.5), Cm(6.5), Cm(0.8), title,
           sz=Pt(18.5), color=C_BLACK, bold=True, font=FONT_B, align=PP_ALIGN.CENTER)
        # 설명
        tb(slide, x, y + Cm(3.5), Cm(6.5), Cm(1), desc,
           sz=Pt(17.5), color=C_TEXT, align=PP_ALIGN.CENTER)

        # 연결 화살표
        if i < 4:
            tb(slide, x + Cm(6.5), y + Cm(0.3), Cm(1), Cm(1), "→",
               sz=Pt(22), color=C_GRAY, align=PP_ALIGN.CENTER)
        x += Cm(7.2)

    # 핵심 메시지
    bottom_note(slide,
        "핵심: 무료 PoC(2~4주)로 실제 병원 환경에서 효과를 직접 확인한 후 도입 여부를 결정할 수 있습니다. "
        "문의부터 결과 리포트까지 전 과정을 SENTIO가 지원합니다.")


# ══════════════════════════════════════════
# 슬라이드 17: PoC 성과 측정 기준 ★★★
# ══════════════════════════════════════════
def s17():
    slide = blank()
    slide_header(slide, "PoC 성과 측정", "명확한 KPI로 객관적 성과 평가")
    slide_num(slide, 17)

    kpis = [
        ("낙상 감지율", "90% 이상",
         "실제 발생 낙상 대비 감지 비율"),
        ("오탐율", "10% 이하",
         "전체 알림 중 실제 낙상이 아닌 비율"),
        ("평균 알림 시간", "5초 이내",
         "낙상 발생 → 알림 도달까지 소요 시간"),
        ("간호사 만족도", "4.0 / 5.0",
         "사용 편의성·알림 정확성·업무 도움 정도"),
    ]

    x = Cm(2)
    for title, target, desc in kpis:
        # 제목
        tb(slide, x, Cm(5), Cm(8), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        # 큰 목표값
        tb(slide, x, Cm(5.8), Cm(8), Cm(1.5), target,
           sz=Pt(36), color=C_RED, bold=True, font=FONT_B)
        # 설명
        tb(slide, x, Cm(7.5), Cm(8), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        x += Cm(8.5)

    bottom_note(slide,
        "PoC 종료 후 상기 KPI 기준 성과 리포트를 제공합니다. 객관적 데이터로 도입 의사결정을 지원합니다.")


# ══════════════════════════════════════════
# 슬라이드 18: 요금제
# ══════════════════════════════════════════
def s18():
    slide = blank()
    slide_header(slide, "요금제", "SaaS 구독형 — 추가 장비 비용 없음")
    slide_num(slide, 18)

    plans = [
        ("Basic", "월 30만원", "카메라 2대",
         ["낙상 감지 · 대시보드 · 알림", "한국어 음성 경고", "이메일 리포트"]),
        ("Pro", "월 70만원", "카메라 8대",
         ["Basic 전체 포함", "다중 카메라 · EMR 연동", "상세 통계 리포트", "우선 기술 지원"]),
        ("Enterprise", "별도 협의", "카메라 무제한",
         ["Pro 전체 포함", "전체 층·건물 커버", "맞춤 개발 · SLA 보장"]),
    ]

    x = Cm(2)
    for name, price, cams, feats in plans:
        # 플랜명
        tb(slide, x, Cm(5), Cm(10), Cm(0.7), name,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        # 가격
        tb(slide, x, Cm(5.7), Cm(10), Cm(1), price,
           sz=Pt(28), color=C_RED if name == "Basic" else C_TEXT, bold=True, font=FONT_B)
        # 카메라
        tb(slide, x, Cm(6.8), Cm(10), Cm(0.5), cams,
           sz=SZ_BODY, color=C_TEXT)
        # 기능
        fy = Cm(7.5)
        for f in feats:
            tb(slide, x + Cm(0.3), fy, Cm(9), Cm(0.5), f"·  {f}",
               sz=SZ_SMALL, color=C_TEXT)
            fy += Cm(0.55)
        x += Cm(11.5)

    # 비용 비교
    bottom_note(slide,
        "초기 장비 비용 0원 (기존 CCTV 활용)  |  SENTIO 30만원~ vs 기존 센서 솔루션 1,000만원+ 초기비용 + 50만원+ 월비용")


# ══════════════════════════════════════════
# 슬라이드 19: 기대 효과 & ROI ★★★
# ══════════════════════════════════════════
def s19():
    slide = blank()
    slide_header(slide, "기대 효과", "도입 시 어떤 효과를 기대할 수 있나?")
    slide_num(slide, 19)

    effects = [
        ("골든타임 단축", "평균 발견 시간 30분 → 30초 이내 (60배 단축)"),
        ("사고 비용 절감", "낙상 1건당 평균 비용 3,000만원 → 사전 예방"),
        ("법적 리스크 감소", "AI 감시 기록으로 주의의무 입증 가능"),
        ("간호 업무 효율", "CCTV 감시 부담 감소 → 핵심 간호 업무에 집중"),
    ]
    y = Cm(5.5)
    for title, desc in effects:
        dot(slide, Cm(2.2), y + Cm(0.15), C_RED)
        tb(slide, Cm(3), y, Cm(28), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, Cm(3), y + Cm(0.9), Cm(28), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        y += Cm(2.5)

    # ROI
    y += Cm(1)
    rect(slide, Cm(2), y, Emu(SW - Cm(4)), Pt(0.5), C_BORDER)
    tb(slide, Cm(2.5), y + Cm(0.3), Cm(28), Cm(0.7), "ROI 시뮬레이션",
       sz=SZ_SUBTITLE, color=C_BLACK, bold=True, font=FONT_B)
    tb(slide, Cm(2.5), y + Cm(1.1), Cm(28), Cm(0.7),
       "월 30만원 투자 → 낙상 사고 1건 예방 시 약 3,000만원 절감  |  투자 대비 100배 효과",
       sz=SZ_BODY, color=C_RED, bold=True, font=FONT_B)


# ══════════════════════════════════════════
# 슬라이드 20: 시장 기회 & 로드맵
# ══════════════════════════════════════════
def s20():
    slide = blank()
    # 우측 이미지 - 너비 기준 비율 유지 (1:1 이미지)
    img(slide, os.path.join(IMG_DIR, "slide8_market.png"),
        Emu(SW // 2 + Cm(1)), 0, width=Emu(SW // 2 - Cm(1)))

    slide_header(slide, "시장 기회", "국내 연 10조 요양 시장, 연 8% 성장 중")
    slide_num(slide, 20)

    markets = [
        ("1차: 요양병원", "전국 1,900개+ — 환자안전관리 의무 대상"),
        ("1차: 요양시설", "전국 5,800개+ — 고령화로 매년 증가"),
        ("2차: 재활·종합병원", "낙상 고위험 병동"),
    ]
    y = Cm(6)
    for title, desc in markets:
        dot(slide, Cm(2.2), y + Cm(0.15), C_RED)
        tb(slide, Cm(3), y, Cm(14), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, Cm(3), y + Cm(0.7), Cm(14), Cm(1.2), desc,
           sz=SZ_BODY, color=C_TEXT)
        y += Cm(2)

    # 로드맵
    y += Cm(0.5)
    tb(slide, Cm(2.2), y, Cm(14), Cm(0.7), "로드맵",
       sz=SZ_SUBTITLE, color=C_BLACK, bold=True, font=FONT_B)
    phases = [
        ("1  단기", "정확도 향상 — 현장 피드백 반영, AI 지속 개선"),
        ("2  중기", "파일럿 실증 — 요양병원 3곳, 2~4주 실증 운영"),
        ("3  장기", "SaaS 정식 출시 — 전국 확대 + 독거노인 모니터링"),
    ]
    ry = y + Cm(0.8)
    for num, desc in phases:
        tb(slide, Cm(2.5), ry, Cm(1.8), Cm(0.6), num,
           sz=Pt(20), color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, Cm(5), ry, Cm(12), Cm(0.6), desc,
           sz=SZ_BODY, color=C_TEXT)
        ry += Cm(0.8)


# ══════════════════════════════════════════
# 슬라이드 21 (마무리): 핵심 요약 + CTA
# ══════════════════════════════════════════
def s21():
    slide = blank()
    slide_header(slide, "핵심 요약", "왜 SENTIO인가?")
    slide_num(slide, 21)

    summaries = [
        ("검증된 정확도", "424개 실제 영상, F1 94.9, 감지율 95.2%"),
        ("즉시 도입 가능", "기존 CCTV 활용, 추가 장비 불필요, 설치 1일"),
        ("프라이버시 보호", "관절점만 분석, 원본 영상 미저장, 개인정보보호법 준수"),
        ("합리적 비용", "초기 비용 0원, 월 30만원부터 — 기존 대비 1/10 수준"),
    ]
    y = Cm(5.5)
    for title, desc in summaries:
        dot(slide, Cm(2.2), y + Cm(0.15), C_RED)
        tb(slide, Cm(3), y, Cm(8), Cm(0.7), title,
           sz=SZ_SUBTITLE, color=C_TEXT, bold=True, font=FONT_B)
        tb(slide, Cm(12), y, Cm(20), Cm(0.7), desc,
           sz=SZ_BODY, color=C_TEXT)
        y += Cm(1.8)

    # CTA
    y += Cm(1.5)
    rect(slide, Cm(2), y, Emu(SW - Cm(4)), Pt(0.5), C_BORDER)
    y += Cm(0.5)
    tb(slide, Cm(2), y, Emu(SW - Cm(4)), Cm(1),
       "지금 무료 PoC를 신청하세요",
       sz=Pt(28), color=C_RED, bold=True, font=FONT_B, align=PP_ALIGN.CENTER)

    y += Cm(1.5)
    tb(slide, Cm(2), y, Emu(SW - Cm(4)), Cm(0.7),
       "1. 전화 또는 이메일 문의  →  2. 병원 환경 원격 확인 (1일)  →  3. 설치 및 PoC 시작",
       sz=SZ_BODY, color=C_TEXT, align=PP_ALIGN.CENTER)

    y += Cm(1.5)
    tb(slide, Cm(2), y, Emu(SW - Cm(4)), Cm(0.5),
       "문의: contact@sentio-ai.com  |  Tel: 000-0000-0000  |  www.sentio-ai.com",
       sz=SZ_SMALL, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# 실행
# ══════════════════════════════════════════
slides = [
    (s01, "표지"),
    (s02, "Executive Summary"),
    (s03, "문제 인식"),
    (s04, "현장의 어려움"),
    (s05, "솔루션 개요"),
    (s06, "경쟁력 비교"),
    (s07, "작동 원리"),
    (s08, "AI 분석 5단계"),
    (s09, "3단계 알림 체계"),
    (s10, "프라이버시 보호"),
    (s11, "검증된 성능"),
    (s12, "유형별 감지율"),
    (s13, "신뢰성 근거"),
    (s14, "주요 기능"),
    (s15, "PoC 제안 개요"),
    (s16, "PoC 진행 절차"),
    (s17, "PoC 성과 측정"),
    (s18, "요금제"),
    (s19, "기대 효과 & ROI"),
    (s20, "시장 기회 & 로드맵"),
    (s21, "마무리 & CTA"),
]

print("SENTIO PoC 제안서 v2 생성 시작...")
for i, (fn, name) in enumerate(slides):
    fn()
    print(f"  [{i+1}/{len(slides)}] {name}")

prs.save(OUTPUT)
print(f"\n[OK] 저장 완료: {OUTPUT}")
print(f"   슬라이드: {len(prs.slides)}장")
print(f"   파일 크기: {os.path.getsize(OUTPUT) / 1024 / 1024:.1f} MB")

# OneDrive에도 복사
import shutil
try:
    shutil.copy2(OUTPUT, OUTPUT_ONEDRIVE)
    print(f"   OneDrive 복사: {OUTPUT_ONEDRIVE}")
except Exception as e:
    print(f"   OneDrive 복사 실패: {e}")
